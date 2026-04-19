import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(18, 52, 86)
BLUE = RGBColor(56, 134, 203)
TEAL = RGBColor(39, 122, 132)
LIGHT = RGBColor(242, 247, 250)
MID = RGBColor(220, 232, 240)
DARK = RGBColor(36, 45, 56)
GREEN = RGBColor(56, 161, 105)
AMBER = RGBColor(221, 107, 32)
GREY = RGBColor(113, 128, 150)
RED = RGBColor(197, 48, 48)
WHITE = RGBColor(255, 255, 255)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(BASE_DIR, "save_tears_miniprogram", "src", "static", "logo.png")
BG_LOGIN_PATH = os.path.join(BASE_DIR, "save_tears_miniprogram", "src", "static", "images", "bg_login.jpg.jpg")
BG_ADMIN_PATH = os.path.join(BASE_DIR, "save_tears_miniprogram", "src", "static", "images", "bg_admin.jpg.jpg")
AVATAR_PATH = os.path.join(BASE_DIR, "save_tears_miniprogram", "src", "static", "images", "logo_avatar.jpg.jpg")


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide, x=12.15, y=0.34, size=0.46):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(x), Inches(y), width=Inches(size), height=Inches(size))


def add_top_band(slide, color=BLUE):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()


def add_side_glow(slide, color=MID):
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.38), Inches(0.65), Inches(0.08), Inches(5.9))
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.fill.background()


def add_title(slide, title, subtitle=None):
    add_top_band(slide)
    add_side_glow(slide)
    add_logo(slide)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(1.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(11)
        r2.font.color.rgb = GREY
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.28), Inches(2.0), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_bullets(slide, items, x=0.9, y=1.6, w=6.0, h=4.9, font_size=20, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = item
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
    return box


def add_small_box(slide, title, items, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for item in items:
        p = tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = item
        r.font.name = "Aptos"
        r.font.size = Pt(11)
        r.font.color.rgb = WHITE


def add_footer(slide, text="Save Tears Team Handover"):
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.72), Inches(7.01), Inches(0.11), Inches(0.11))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BLUE
    dot.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.02), Inches(12.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = GREY


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_top_band(slide)
    if os.path.exists(BG_LOGIN_PATH):
        slide.shapes.add_picture(BG_LOGIN_PATH, Inches(7.95), Inches(0.28), width=Inches(5.38), height=Inches(7.22))
        overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.95), Inches(0.28), Inches(5.38), Inches(7.22))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = NAVY
        overlay.fill.transparency = 0.34
        overlay.line.fill.background()
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.85), Inches(6.7), Inches(5.6))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = MID
    panel.shadow.inherit = False
    add_logo(slide, x=0.92, y=1.02, size=0.62)
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(8.5), Inches(1.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Save Tears"
    r.font.name = "Aptos Display"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Water Conservation Management System"
    r2.font.name = "Aptos"
    r2.font.size = Pt(18)
    r2.font.color.rgb = TEAL
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(2.0), Inches(2.0), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BLUE
    badge.line.color.rgb = BLUE
    badge_tf = badge.text_frame
    badge_tf.clear()
    bp = badge_tf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "TEAM HANDOVER"
    br.font.name = "Aptos"
    br.font.size = Pt(12)
    br.font.bold = True
    br.font.color.rgb = WHITE
    add_bullets(
        slide,
        [
            "Mini-program frontend for end users",
            "FastAPI backend for user and data management",
            "Focused on registration, login, and room-based water data lookup",
            "Current status: the minimum working flow is complete",
        ],
        x=0.95,
        y=2.55,
        w=6.8,
        h=2.8,
        font_size=18,
    )
    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(1.1), Inches(3.95), Inches(4.9))
    side.fill.solid()
    side.fill.fore_color.rgb = WHITE
    side.fill.transparency = 0.08
    side.line.color.rgb = WHITE
    tf2 = side.text_frame
    tf2.clear()
    if os.path.exists(AVATAR_PATH):
        slide.shapes.add_picture(AVATAR_PATH, Inches(9.72), Inches(1.38), width=Inches(1.25), height=Inches(1.23))
    for idx, line in enumerate(["Audience: UK teammates", "Format: Team handover", "Length: 20-25 mins", "Includes: short live demo"]):
        p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        if idx == 0:
            p.space_before = Pt(92)
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(16)
        r.font.color.rgb = NAVY
    add_footer(slide)


def purpose_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Purpose of This Handover")
    add_small_box(slide, "Working Core", ["Show what is already working"], 0.9, 1.8, 2.7, 1.5, BLUE)
    add_small_box(slide, "Scope", ["Separate real logic from UI prototypes"], 3.9, 1.8, 2.7, 1.5, TEAL)
    add_small_box(slide, "Structure", ["Explain how the codebase is organized"], 6.9, 1.8, 2.7, 1.5, NAVY)
    add_small_box(slide, "Next Phase", ["Propose a practical team split"], 9.9, 1.8, 2.5, 1.5, GREEN)
    add_bullets(
        slide,
        [
            "This is a handover and onboarding session, not a final defense presentation.",
            "The goal is to help teammates understand the current implementation and continue development efficiently.",
        ],
        x=1.0,
        y=4.1,
        w=11.2,
        h=1.8,
        font_size=20,
    )
    add_footer(slide)


def architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_title(slide, "Current System Architecture")
    shapes = [
        (1.0, 2.2, 2.2, 1.0, "Mini-program\npages", BLUE),
        (4.0, 2.2, 2.2, 1.0, "src/api/\nindex.ts", TEAL),
        (7.0, 2.2, 2.2, 1.0, "FastAPI\nroutes", NAVY),
        (10.0, 2.2, 2.2, 1.0, "SQLite / MySQL\ndatabase", GREEN),
    ]
    for x, y, w, h, text, color in shapes:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = "Aptos"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE
    for x in [3.2, 6.2, 9.2]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.47), Inches(0.5), Inches(0.46))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GREY
        arrow.line.color.rgb = GREY
    add_bullets(
        slide,
        [
            "Page interaction -> API wrapper -> backend route -> database -> JSON response -> page rendering",
            "Frontend: save_tears_miniprogram",
            "Backend entry points: save_tears_backend/api.py and save_tears_backend/main.py",
        ],
        x=1.0,
        y=4.1,
        w=11.0,
        h=1.8,
        font_size=18,
    )
    add_footer(slide)


def working_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "What Is Already Working")
    steps = [
        ("Register", BLUE),
        ("Log in", TEAL),
        ("Store user\nlocally", NAVY),
        ("Query by\nroom number", GREEN),
        ("Show three\nroom-based pages", AMBER),
    ]
    x = 0.8
    for idx, (label, color) in enumerate(steps):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.1), Inches(2.0), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 2.05), Inches(2.37), Inches(0.45), Inches(0.45))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREY
            arrow.line.color.rgb = GREY
        x += 2.45
    add_bullets(
        slide,
        [
            "Working room-based lookup pages: water flow, sewage turbidity, and water bill.",
            "This flow is backend-connected and persists data through the database.",
        ],
        x=1.0,
        y=4.15,
        w=11.0,
        h=1.6,
        font_size=19,
    )
    add_footer(slide)


def frontend_status_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_title(slide, "Frontend Status by Page")
    add_small_box(slide, "Working business pages", ["login", "register", "water-flow", "sewage-turbidity", "water-bill"], 0.9, 1.8, 3.6, 3.1, GREEN)
    add_small_box(slide, "Partially dynamic pages", ["home", "profile"], 4.9, 1.8, 3.0, 3.1, AMBER)
    add_small_box(slide, "Prototype / mock page", ["admin"], 8.3, 1.8, 3.9, 3.1, GREY)
    add_bullets(
        slide,
        [
            "Some pages are API-driven and already usable.",
            "Others are visually complete enough to show direction, but still need real business logic.",
        ],
        x=1.0,
        y=5.3,
        w=11.0,
        h=1.2,
        font_size=18,
    )
    add_footer(slide)


def backend_api_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Main Backend Interfaces")
    add_small_box(
        slide,
        "Used by the mini-program now",
        ["POST /register", "POST /login", "GET /water_flow/{room_number}", "GET /sewage_turbidity/{room_number}", "GET /water_bill/{room_number}"],
        0.9,
        1.8,
        5.3,
        3.6,
        BLUE,
    )
    add_small_box(
        slide,
        "Implemented but not connected in UI",
        ["POST /water_flow", "POST /sewage_turbidity", "POST /water_bill", "GET /users"],
        6.7,
        1.8,
        5.6,
        3.6,
        TEAL,
    )
    add_footer(slide)


def codebase_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_title(slide, "How to Read the Codebase Quickly")
    add_small_box(slide, "1. pages.json", ["Routes", "Page list", "Tab bar"], 0.9, 2.0, 3.6, 2.4, BLUE)
    add_small_box(slide, "2. src/api/index.ts", ["Request wrapper", "API functions"], 4.9, 2.0, 3.6, 2.4, TEAL)
    add_small_box(slide, "3. save_tears_backend/api.py", ["Models", "Schemas", "Routes", "Database queries"], 8.9, 2.0, 3.4, 2.4, NAVY)
    add_bullets(
        slide,
        ["These three files give the fastest complete overview for a new teammate."],
        x=1.0,
        y=5.1,
        w=11.0,
        h=0.9,
        font_size=20,
    )
    add_footer(slide)


def gaps_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Current Gaps and Risks")
    left_items = [
        "Home page content is mostly static",
        "Profile feature buttons are not connected",
        "Admin page uses mock data only",
    ]
    right_items = [
        "Registration UI collects email, but backend does not store it",
        "Authentication is local-storage based, not token-based",
        "Passwords are currently compared in plaintext",
    ]
    add_small_box(slide, "Product and UI gaps", left_items, 0.9, 1.8, 5.4, 3.6, AMBER)
    add_small_box(slide, "Technical and security gaps", right_items, 6.7, 1.8, 5.6, 3.6, RED)
    add_footer(slide)


def team_split_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_title(slide, "Suggested Team Ownership")
    add_small_box(slide, "Frontend feature completion", ["home page", "profile page", "admin page", "charting"], 0.8, 1.8, 2.9, 3.2, BLUE)
    add_small_box(slide, "Backend and security", ["email field", "password hashing", "auth model", "API cleanup"], 3.95, 1.8, 2.9, 3.2, TEAL)
    add_small_box(slide, "Data input and visualization", ["connect submit-data APIs", "add data entry forms", "build real chart pages"], 7.1, 1.8, 2.9, 3.2, NAVY)
    add_small_box(slide, "Testing and integration", ["environment setup", "sample data", "regression checks", "handover docs"], 10.25, 1.8, 2.2, 3.2, GREEN)
    add_footer(slide)


def demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Short Demo Flow")
    add_bullets(
        slide,
        [
            "1. Register a test user",
            "2. Log in",
            "3. Show the username on the home page",
            "4. Open the three room-based data pages",
            "5. Confirm that records are loaded from backend data",
        ],
        x=1.0,
        y=1.9,
        w=6.7,
        h=3.4,
        font_size=21,
    )
    add_small_box(slide, "Before the meeting", ["backend running", "sample records ready", "local storage cleaned"], 8.1, 2.1, 3.8, 2.3, TEAL)
    add_footer(slide)


def sprint_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_title(slide, "Next Sprint Proposal")
    add_small_box(slide, "Stage 1", ["stabilize the backend model", "add authentication and security basics"], 1.0, 2.0, 3.5, 2.2, BLUE)
    add_small_box(slide, "Stage 2", ["turn placeholder pages into real features", "connect data input and charting"], 4.9, 2.0, 3.5, 2.2, TEAL)
    add_small_box(slide, "Stage 3", ["testing", "polish", "deployment and final presentation readiness"], 8.8, 2.0, 3.5, 2.2, GREEN)
    add_bullets(
        slide,
        ["The project already has a solid working core. The next step is to make it more complete and team-owned."],
        x=1.0,
        y=5.0,
        w=11.0,
        h=1.0,
        font_size=20,
    )
    add_footer(slide)


def build_presentation(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    title_slide(prs)
    purpose_slide(prs)
    architecture_slide(prs)
    working_flow_slide(prs)
    frontend_status_slide(prs)
    backend_api_slide(prs)
    codebase_slide(prs)
    gaps_slide(prs)
    team_split_slide(prs)
    demo_slide(prs)
    sprint_slide(prs)
    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("Save_Tears_Team_Handover.pptx")
    print("Generated Save_Tears_Team_Handover.pptx")
